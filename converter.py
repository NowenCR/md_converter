"""
converter.py — Core conversion engine for MD Forge
Supports: PDF, DOCX, DOC, TXT, RTF, HTML, HTM, CSV, ODT, EPUB, PPTX, PPT, MD
"""

import os
import re
import csv
from pathlib import Path


class FileConverter:
    SUPPORTED = {
        ".pdf": "_convert_pdf",
        ".docx": "_convert_docx",
        ".doc": "_convert_doc",
        ".txt": "_convert_txt",
        ".rtf": "_convert_rtf",
        ".html": "_convert_html",
        ".htm": "_convert_html",
        ".csv": "_convert_csv",
        ".odt": "_convert_odt",
        ".epub": "_convert_epub",
        ".pptx": "_convert_pptx",
        ".ppt": "_convert_ppt",
        ".md": "_convert_md",
    }

    def convert(self, filepath: str) -> str:
        ext = Path(filepath).suffix.lower()
        method_name = self.SUPPORTED.get(ext)
        if not method_name:
            raise ValueError(f"Unsupported file type: '{ext}'")
        method = getattr(self, method_name)
        return method(filepath)

    # ─── PDF ────────────────────────────────────────────────────────────────

    def _convert_pdf(self, filepath: str) -> str:
        try:
            import pdfplumber
            lines = []
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    lines.append(f"\n---\n*Page {page_num}*\n")
                    # Extract tables first
                    tables = page.extract_tables()
                    table_bboxes = [t.bbox for t in page.find_tables()] if hasattr(page, "find_tables") else []

                    if tables:
                        for table in tables:
                            if table and len(table) > 0:
                                lines.append(self._table_to_md(table))

                    # Extract text
                    text = page.extract_text()
                    if text:
                        lines.append(self._clean_text_to_md(text))
            return "\n".join(lines)
        except ImportError:
            # Fallback to PyPDF2
            try:
                import PyPDF2
                lines = []
                with open(filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages, 1):
                        lines.append(f"\n---\n*Page {i}*\n")
                        text = page.extract_text() or ""
                        lines.append(self._clean_text_to_md(text))
                return "\n".join(lines)
            except ImportError:
                raise ImportError(
                    "No PDF library found. Install with:\n  pip install pdfplumber\nor\n  pip install PyPDF2"
                )

    # ─── DOCX ───────────────────────────────────────────────────────────────

    def _convert_docx(self, filepath: str) -> str:
        try:
            import docx
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        doc = docx.Document(filepath)
        lines = []
        title = Path(filepath).stem
        lines.append(f"# {title}\n")

        for element in doc.element.body:
            tag = element.tag.split("}")[-1]

            if tag == "p":
                para = docx.text.paragraph.Paragraph(element, doc)
                text = para.text.strip()
                if not text:
                    lines.append("")
                    continue
                style = para.style.name.lower()
                if "heading 1" in style:
                    lines.append(f"# {text}")
                elif "heading 2" in style:
                    lines.append(f"## {text}")
                elif "heading 3" in style:
                    lines.append(f"### {text}")
                elif "heading 4" in style:
                    lines.append(f"#### {text}")
                elif "list bullet" in style or "list paragraph" in style:
                    lines.append(f"- {text}")
                elif "list number" in style:
                    lines.append(f"1. {text}")
                else:
                    # Inline formatting
                    md_line = self._docx_para_to_md(para)
                    lines.append(md_line)

            elif tag == "tbl":
                from docx.table import Table
                table = Table(element, doc)
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                lines.append(self._table_to_md(rows))

        return "\n".join(lines)

    def _docx_para_to_md(self, para) -> str:
        import docx
        parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            if run.bold and run.italic:
                text = f"***{text}***"
            elif run.bold:
                text = f"**{text}**"
            elif run.italic:
                text = f"*{text}*"
            if run.underline:
                text = f"<u>{text}</u>"
            parts.append(text)
        return "".join(parts)

    # ─── DOC (legacy) ───────────────────────────────────────────────────────

    def _convert_doc(self, filepath: str) -> str:
        # Try antiword or LibreOffice as system tools
        import subprocess
        try:
            result = subprocess.run(
                ["antiword", filepath], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                return self._clean_text_to_md(result.stdout)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        try:
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "txt", "--outdir", "/tmp", filepath],
                capture_output=True, timeout=60
            )
            tmp = f"/tmp/{Path(filepath).stem}.txt"
            if os.path.exists(tmp):
                with open(tmp, "r", encoding="utf-8", errors="replace") as f:
                    return self._clean_text_to_md(f.read())
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        raise RuntimeError(
            "Cannot convert .doc files without antiword or LibreOffice installed.\n"
            "Install antiword: sudo apt install antiword\n"
            "Or save as .docx from Word."
        )

    # ─── TXT ────────────────────────────────────────────────────────────────

    def _convert_txt(self, filepath: str) -> str:
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        for enc in encodings:
            try:
                with open(filepath, "r", encoding=enc) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()

        title = Path(filepath).stem
        lines = [f"# {title}\n"]
        lines.append(self._clean_text_to_md(content))
        return "\n".join(lines)

    # ─── RTF ────────────────────────────────────────────────────────────────

    def _convert_rtf(self, filepath: str) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                rtf = f.read()
            text = rtf_to_text(rtf)
            title = Path(filepath).stem
            return f"# {title}\n\n" + self._clean_text_to_md(text)
        except ImportError:
            raise ImportError("striprtf not installed. Run: pip install striprtf")

    # ─── HTML ───────────────────────────────────────────────────────────────

    def _convert_html(self, filepath: str) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")

        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")

        # Remove scripts and styles
        for tag in soup(["script", "style", "head", "nav", "footer"]):
            tag.decompose()

        lines = []
        for element in soup.find_all(True):
            tag = element.name
            text = element.get_text(strip=True)
            if not text:
                continue
            if tag == "h1":
                lines.append(f"# {text}\n")
            elif tag == "h2":
                lines.append(f"## {text}\n")
            elif tag == "h3":
                lines.append(f"### {text}\n")
            elif tag == "h4":
                lines.append(f"#### {text}\n")
            elif tag == "h5":
                lines.append(f"##### {text}\n")
            elif tag == "h6":
                lines.append(f"###### {text}\n")
            elif tag == "p":
                lines.append(f"{text}\n")
            elif tag == "li":
                lines.append(f"- {text}")
            elif tag == "blockquote":
                lines.append(f"> {text}\n")
            elif tag == "code":
                lines.append(f"`{text}`")
            elif tag == "pre":
                lines.append(f"```\n{text}\n```\n")
            elif tag == "a":
                href = element.get("href", "#")
                lines.append(f"[{text}]({href})")
            elif tag == "img":
                alt = element.get("alt", "image")
                src = element.get("src", "")
                lines.append(f"![{alt}]({src})")
            elif tag == "table":
                rows = []
                for tr in element.find_all("tr"):
                    row = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                    rows.append(row)
                lines.append(self._table_to_md(rows))

        # Deduplicate consecutive identical lines
        result = []
        prev = None
        for line in lines:
            if line != prev:
                result.append(line)
            prev = line

        return "\n".join(result)

    # ─── CSV ────────────────────────────────────────────────────────────────

    def _convert_csv(self, filepath: str) -> str:
        title = Path(filepath).stem
        lines = [f"# {title}\n"]
        rows = []
        with open(filepath, "r", encoding="utf-8-sig", errors="replace", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        if rows:
            lines.append(self._table_to_md(rows))
        lines.append(f"\n*{len(rows)} rows total*")
        return "\n".join(lines)

    # ─── ODT ────────────────────────────────────────────────────────────────

    def _convert_odt(self, filepath: str) -> str:
        try:
            from odf.opendocument import load
            from odf.text import P, H, List, ListItem
            from odf import teletype
        except ImportError:
            raise ImportError("odfpy not installed. Run: pip install odfpy")

        doc = load(filepath)
        title = Path(filepath).stem
        lines = [f"# {title}\n"]

        for elem in doc.text.childNodes:
            tag = elem.__class__.__name__
            text = teletype.extractText(elem).strip()
            if not text:
                continue
            if tag == "H":
                level = int(elem.getAttribute("outlinelevel") or 1)
                prefix = "#" * min(level, 6)
                lines.append(f"{prefix} {text}")
            elif tag == "P":
                lines.append(text)
            elif tag in ("List", "ListItem"):
                lines.append(f"- {text}")

        return "\n".join(lines)

    # ─── EPUB ───────────────────────────────────────────────────────────────

    def _convert_epub(self, filepath: str) -> str:
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError(
                "ebooklib/beautifulsoup4 not installed.\nRun: pip install ebooklib beautifulsoup4"
            )

        book = epub.read_epub(filepath)
        lines = []
        title = book.get_metadata("DC", "title")
        if title:
            lines.append(f"# {title[0][0]}\n")

        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            lines.append(self._clean_text_to_md(text))
            lines.append("\n---\n")

        return "\n".join(lines)

    # ─── PPTX ───────────────────────────────────────────────────────────────

    def _convert_pptx(self, filepath: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(filepath)
        title = Path(filepath).stem
        lines = [f"# {title}\n"]

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title = self._pptx_slide_title(slide)
            if slide_title:
                lines.append(f"\n---\n\n## Slide {slide_num}: {slide_title}\n")
            else:
                lines.append(f"\n---\n\n## Slide {slide_num}\n")

            # Process shapes in reading order (top-to-bottom, left-to-right)
            shapes = sorted(
                slide.shapes,
                key=lambda s: (s.top if s.top is not None else 0,
                               s.left if s.left is not None else 0)
            )

            for shape in shapes:
                # Skip the title shape (already handled above)
                if shape == self._get_title_shape(slide):
                    continue

                # Tables
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text_frame.text.strip() for cell in row.cells])
                    lines.append(self._table_to_md(rows))
                    continue

                # Text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        level = para.level  # 0 = top level, 1+ = indented
                        # Detect if it looks like a bullet list item
                        is_bullet = self._is_bullet_para(para)
                        if is_bullet:
                            indent = "  " * level
                            lines.append(f"{indent}- {text}")
                        else:
                            # Body text: treat level-0 as paragraph, deeper as bullet
                            if level == 0:
                                lines.append(f"\n{text}\n")
                            else:
                                indent = "  " * (level - 1)
                                lines.append(f"{indent}- {text}")

            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""
                if notes_text:
                    lines.append(f"\n> **Notes:** {notes_text}\n")

        return "\n".join(lines)

    def _pptx_slide_title(self, slide) -> str:
        """Extract the title text from a slide, or empty string."""
        shape = self._get_title_shape(slide)
        if shape and shape.has_text_frame:
            return shape.text_frame.text.strip()
        return ""

    def _get_title_shape(self, slide):
        """Return the title placeholder shape if present."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                # idx 0 = title, idx 1 = body/subtitle on title slide
                if ph and ph.idx == 0:
                    return shape
        return None

    def _is_bullet_para(self, para) -> bool:
        """Heuristic: treat as bullet if level > 0 or explicit bullet char present."""
        if para.level > 0:
            return True
        text = para.text.strip()
        if text.startswith(("•", "–", "—", "-", "▪", "◦", "○", "●")):
            return True
        # Check XML for bullet formatting
        try:
            from lxml import etree
            pPr = para._p.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
            )
            if pPr is not None:
                buNone = pPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
                )
                if buNone is not None:
                    return False
                buChar = pPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
                )
                buAutoNum = pPr.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum"
                )
                if buChar is not None or buAutoNum is not None:
                    return True
        except Exception:
            pass
        return False

    # ─── PPT (legacy binary) ────────────────────────────────────────────────

    def _convert_ppt(self, filepath: str) -> str:
        """
        Legacy .ppt (PowerPoint 97-2003) conversion.
        Strategy 1: LibreOffice headless → .pptx → reuse _convert_pptx
        Strategy 2: LibreOffice headless → .txt → clean_text_to_md
        Strategy 3: python-pptx direct attempt (sometimes works for newer .ppt)
        """
        import subprocess
        import tempfile

        tmp_dir = tempfile.mkdtemp()

        # Strategy 1: convert via LibreOffice to pptx
        try:
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pptx",
                 "--outdir", tmp_dir, filepath],
                capture_output=True, timeout=60
            )
            pptx_path = os.path.join(tmp_dir, Path(filepath).stem + ".pptx")
            if os.path.exists(pptx_path):
                return self._convert_pptx(pptx_path)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Strategy 2: convert via LibreOffice to txt
        try:
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "txt",
                 "--outdir", tmp_dir, filepath],
                capture_output=True, timeout=60
            )
            txt_path = os.path.join(tmp_dir, Path(filepath).stem + ".txt")
            if os.path.exists(txt_path):
                with open(txt_path, "r", encoding="utf-8", errors="replace") as f:
                    return self._clean_text_to_md(f.read())
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

        # Strategy 3: try python-pptx directly (may work for some .ppt files)
        try:
            return self._convert_pptx(filepath)
        except Exception:
            pass

        raise RuntimeError(
            "Cannot convert .ppt (legacy) files without LibreOffice.\n"
            "Install: https://www.libreoffice.org/download/download/\n"
            "Or re-save the file as .pptx from PowerPoint."
        )

    # ─── MD (passthrough) ───────────────────────────────────────────────────

    def _convert_md(self, filepath: str) -> str:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # ─── Helpers ────────────────────────────────────────────────────────────

    def _table_to_md(self, rows: list) -> str:
        if not rows:
            return ""
        # Normalize row lengths
        max_cols = max(len(r) for r in rows)
        rows = [r + [""] * (max_cols - len(r)) for r in rows]

        def sanitize(cell):
            return str(cell).replace("|", "\\|").replace("\n", " ").strip()

        lines = []
        header = rows[0]
        lines.append("| " + " | ".join(sanitize(c) for c in header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(sanitize(c) for c in row) + " |")
        return "\n".join(lines) + "\n"

    def _clean_text_to_md(self, text: str) -> str:
        if not text:
            return ""
        lines = text.splitlines()
        md_lines = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                md_lines.append("")
                continue
            # Detect all-caps short lines as potential headings
            if stripped.isupper() and len(stripped) < 80 and len(stripped.split()) <= 8:
                md_lines.append(f"## {stripped.title()}")
            else:
                md_lines.append(stripped)
        # Collapse multiple blank lines
        result = re.sub(r"\n{3,}", "\n\n", "\n".join(md_lines))
        return result
